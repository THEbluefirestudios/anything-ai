from huggingface_hub import InferenceClient
import os
import platform
import subprocess
from time import sleep
import threading as thread #import block for basic features
import textwrap
import msvcrt

import AppOpener as app# import block for opening apps and web
import webbrowser as web

import json
from pptx import Presentation
from pptx.util import Pt
from docx import Document # import block for for file gen 
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
import csv
import openpyxl
import requests
from io import BytesIO

from PIL import Image
import base64
import tkinter as tk
from tkinter import filedialog # import block for multimodal input
from pathlib import Path

import datetime
from random import choice # import block for the tui
from colorama import Fore as col



#THIS THING DOES ALMOST EVERYTHING MAINSTREAM CLOSED SOURCE AI DOES: LIKE MULTIMODAL INPUT, FILE GENERATION, CHAT MEMORY, .ETC

hf_token_1 = str(os.environ.get("HF_ACCESS_TOKEN"))# adding this block for the non savvy 
hf_token_2 = str(os.environ.get("HF_ACCESS_TOKEN_2"))
hf_token_3 = str(os.environ.get("HF_ACCESS_TOKEN_3"))

if hf_token_1 == "None" and hf_token_2 == "None" and hf_token_3 == "None":
    print(col.RED + "\n[Anything AI] No Hugging Face tokens found.")
    print(col.YELLOW + "Please set at least one of the following environment variables:")
    print(col.WHITE + "  HF_ACCESS_TOKEN")
    print(col.WHITE + "  HF_ACCESS_TOKEN_2")
    print(col.WHITE + "  HF_ACCESS_TOKEN_3")
    print(col.LIGHTBLACK_EX + "\nGet a free token at: https://huggingface.co/settings/tokens")
    input("\nPress Enter to exit...")
    exit()

import os
import json


# look mom i made my own input function!


def read_input(prompt, seed=None):
    global attachfile, filebar_files
    if seed is None:
        print(prompt, end='', flush=True)
        buffer = ""
    else:
        print(prompt + seed, end='', flush=True) 
        buffer = seed

    while True:
        ch = msvcrt.getwch()
        if ch == '\r':
            print()
            return buffer
        elif ch == '\x08':
            if buffer:
                buffer = buffer[:-1]
                print('\b \b', end='', flush=True)
        elif ch == '\xe0' or ch == '\x00':
            special = msvcrt.getwch()
            if special == 'R':
                picked = pick_file()
                if picked:
                    attachfile.append(picked)
                    ext = picked.split('.')[-1].lower()
                    if ext in ("png", "jpg", "jpeg", "gif", "webp"):
                        icon = "▣"
                    elif ext in ("py", "java", "js", "html", "kt", "c", "cpp", "cs", "css", "pyw", "ipynb"):
                        icon = "<\\>"
                    elif ext in ("txt", "docx", "csv", "pptx", "xlsx"):
                        icon = "₠"
                    else:
                        icon = "▮"
                    filebar_files.append((icon, Path(picked).name))
                    draw_file_bar()
                    print(f"\r{prompt}{buffer}", end='', flush=True)
        elif ch == '\x1b':
            pass
        else:
            buffer += ch
            print(ch, end='', flush=True)



APP_NAME = "Anything AI" # was gonna originally name this 'artificial inteligence AI' XD
hf_token_1 = str(os.environ.get("HF_ACCESS_TOKEN"))
hf_token_2 = str(os.environ.get("HF_ACCESS_TOKEN_2"))
hf_token_3 = str(os.environ.get("HF_ACCESS_TOKEN_3"))

client1 = InferenceClient(token=hf_token_1)
client2 = InferenceClient(token=hf_token_2)# huggign face should give atleast something more than 0.10 dolla per month
client3 = InferenceClient(token=hf_token_3)

size = os.get_terminal_size()

width = size.columns

header = [' █████╗ ███╗   ██╗██╗   ██╗████████╗██╗  ██╗██╗███╗   ██╗ ██████╗      █████╗ ██╗', '██╔══██╗████╗  ██║╚██╗ ██╔╝╚══██╔══╝██║  ██║██║████╗  ██║██╔════╝     ██╔══██╗██║','███████║██╔██╗ ██║ ╚████╔╝    ██║   ███████║██║██╔██╗ ██║██║  ███╗    ███████║██║', '██╔══██║██║╚██╗██║  ╚██╔╝     ██║   ██╔══██║██║██║╚██╗██║██║   ██║    ██╔══██║██║','██║  ██║██║ ╚████║   ██║      ██║   ██║  ██║██║██║ ╚████║╚██████╔╝    ██║  ██║██║', '╚═╝  ╚═╝╚═╝  ╚═══╝   ╚═╝      ╚═╝   ╚═╝  ╚═╝╚═╝╚═╝  ╚═══╝ ╚═════╝     ╚═╝  ╚═╝╚═╝']
# list final boss
context = []
username = str(os.getlogin())# i dont collect data,t his is how i know ur name

def determine_time_of_day(time): # useful
    hour = time.hour
    if 5 <= hour < 12:
        return "morning"
    elif 12 <= hour < 17:
        return "afternoon"
    elif 17 <= hour < 21:
        return "evening"
    else:
        return "night"
    
daytime = determine_time_of_day(datetime.datetime.now())

welcome_messages = [# u just KNOW i was looking at the claude homepage for reference
    f"Good {daytime}, {username}. What are we building today?",
    f"Back at it, {username}? Good {daytime}, let's get into it.",
    f"Hey {username}, good {daytime}. What do you need?",
    f"Good {daytime}, {username}. Pick up where we left off?",
    f"{username}, good {daytime}. What can I help you crack today?",
    f"Hey {username}, what's the plan?",
    f"Back again, {username}?",
    f"Oh, {username}'s here. What are we breaking today?",
    f"Welcome back, {username}. Ready when you are.",
    f"{username}. What do you need?",
    f"Ah, {username}. What's on your mind?",
    f"Good to see you, {username}. Let's get to work.",
    f"There you are, {username}. What are we cracking?",
    f"Good {daytime}. What are we working on?",
    f"Good {daytime}. Ready when you are.",
    f"Good {daytime}. Got something for me?",
    f"Good {daytime}. Let's get into it.",
    f"What are we building today?",
    f"Back at it again?",
    f"Ready when you are.",
    f"What do you need?",
    f"Let's get to work.",
    f"Got a problem for me?",
    f"What's on the agenda?",
    f"What are we breaking today?",
    f"Pick up where we left off?",
    f"Let's make something.",
]

summarized_context = ""
message = choice(welcome_messages)

# parsation of files
def parse_text_file(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def parse_docx_file(path):
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

def parse_csv_file(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def parse_xlsx_file(path):
    wb = openpyxl.load_workbook(path)
    ws = wb.active
    rows = []
    for row in ws.iter_rows(values_only=True):  # type: ignore
        rows.append("\t".join(str(c) if c is not None else "" for c in row))
    return "\n".join(rows)

def parse_pptx_file(path):
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        text = " ".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))  # type: ignore
        slides.append(f"Slide {i+1}: {text}")
    return "\n".join(slides)

def parse_code_file(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

# cool effect u see on ai chatting interfaces
def type_print(text, delay):
    for line in text.splitlines():
        wrapped_lines = textwrap.wrap(line, width=width) if line.strip() else ['']
        for wrapped in wrapped_lines:
            for char in wrapped:
                print(char, end='', flush=True)
                sleep(delay)
            print()
    print()

def ui_header(): # im gonna use this fuc only one time!
    global daytime, header, welcome_messages, username, width, message
    heading_spaces = (width - len(header[1]))/2

    message_spaces = (width - len(message))/2
    if width >= len(header[1]):
        for x in range(len(header)):
            print(col.LIGHTBLUE_EX+(round(heading_spaces)*' ')+header[x])
    else:
        heading_spaces = (width - len('ANYTHING AI'))/2
        print(col.LIGHTBLUE_EX+ (' '* round(heading_spaces)) +'ANYTHING AI')
        
    print('\n'*2)
    print(round(message_spaces)*' '+message)
        
        

def open_file(path): 
    current_os = platform.system()

    if current_os == "Windows":# doing some jugaad to support all oses, but this thing will ONLY ship for windows, im too lazy to remove it now
        os.startfile(path)
    elif current_os == "Darwin":  
        subprocess.run(["open", path])
    elif current_os == "Linux":
        subprocess.run(["xdg-open", path])
    else:
        print("Unsupported operating system")

def run_terminal_command(command): #i DEFINITELY TRUST AI TO RUN COMMANDS FOR ME
    global width
    type_print(col.YELLOW+'AI has requested to run the given command: ', 0.01)
    print(col.LIGHTBLUE_EX+'\n')
    print(width*'-'+'\n')
    print('   '+command+ '\n')
    print(width*'-'+'\n')
    type_print(col.YELLOW+'\nDo you accept or reject the execution of this command? [y/n]'+ col.WHITE, 0.01)
    choice = ''
    while choice.lower() not in ['y', 'n']:
        choice = str(input())
        if choice.lower() not in ['y','n']:
            print(col.YELLOW+'Invalid input, please enter y or n'+ col.WHITE)
    if choice.lower() == 'y':
        print(f"{col.CYAN}Terminal"+ ('-'*(width-8))+ '\n' + col.WHITE)
        subprocess.run(command, shell=True)
        print('\n'+ col.CYAN + ('-'* width)+ '\n')
        context.append('User accepted AI\'s request to run command')
        print(col.LIGHTGREEN_EX+f'\nExecuted command:{command}\n')
    else:
        context.append('User rejected AI\'s request to run command')

def get_image_base64_uri(image_path): #i think base64 therefore i am base64 (stardance reference)
    with Image.open(image_path) as img:
        buffered = BytesIO()
        img_format = img.format if img.format else "JPEG"
        img.save(buffered, format=img_format)
        img_base64 = base64.b64encode(buffered.getvalue()).decode('utf-8')
        return f"data:image/{img_format.lower()};base64,{img_base64}"

def pick_file(): # love to navigate my messy folder structure testing this!
    root = tk.Tk()
    root.withdraw()
    root.attributes('-topmost', True)
    
    path = filedialog.askopenfilename(
        title="Select a File",
        filetypes=[
            ("Image files", "*.png *.jpg *.jpeg *.gif *.webp"),
            ("Text and Documents", "*.txt *.docx *.csv *.pptx *xlsx"),
            ("Code files", "*.py *.java *.js *.html *.kt *.c *.cpp *.cs *.css *.pyw *.ipynb"),
            ("All files", "*.*")
        ]
    )
    
    root.destroy()

    if not path:
        return None
    path = str(path).replace("\\", "\\\\")

    return path

def fetch_wikimedia_image(query):# thats why wikipedia is the GOAT.... THE GOAT
    try:
        search_url = "https://en.wikipedia.org/w/api.php"
        search_params = {
            "action": "query",
            "format": "json",
            "generator": "search",
            "gsrsearch": query,
            "gsrlimit": 1,
            "prop": "pageimages",
            "piprop": "original",
        }
        myheaders = {
            "User-Agent": "Anything AI/1.0 (THEbluefirestudios.github.io)" #very real header
        }
        response = requests.get(search_url, params=search_params, headers=myheaders, timeout=5)
        data = response.json()
        pages = data.get("query", {}).get("pages", {})
        for page in pages.values():
            img_url = page.get("original", {}).get("source")
            if img_url and img_url.lower().endswith((".jpg", ".jpeg", ".png")):
                img_data = requests.get(img_url, headers=myheaders, timeout=5).content
                return BytesIO(img_data)
    except Exception as e:
        print(str(e))
        pass
    return None


def get_output_dir(): # i didnt need to make this function
    username = os.getlogin()
    output_dir = f"C:\\Users\\{username}\\{APP_NAME}\\output\\"
    os.makedirs(output_dir, exist_ok=True)
    return output_dir





def maybe_summarize_context(): #MAYBE is a key word here
    global context, summarized_context
    if len(context) <= 10:
        return
    to_summarize = context[:-10]
    keep = context[-10:]
    raw = "\n".join(to_summarize)
    system_prompt = (
        f"You are a local companion memory engine. Update the summarised context (currently: {summarized_context}) to contain details from the given conversation snippet. Keep it ultra-condensed, preserving only core settings, app errors, preferences, computer code or specific details of the conversation. Make sure to not use ANY markdown formatting in the summary, keep everything in plaintext. Make sure to keep computer code in ITS ENTIRETY, do NOT shorten computer code. Start working immediately."
    )
    try:
        summarized_context = call_with_fallback(
            "meta-llama/Llama-3.1-8B-Instruct",
            [{"role": "system", "content": system_prompt}, {"role": "user", "content": raw}],
            500, 0.1
        )
        context = keep
    except:
        pass


def get_context_string():
    return f"Summarized history:\n{summarized_context}\n\nRecent exchanges:\n{chr(10).join(context)}"

# big AI imitation here
def load_memory():
    global summarized_context
    memory_path = f"C:\\Users\\{os.getlogin()}\\Anything AI\\memory\\memory.txt"
    if os.path.exists(memory_path):
        with open(memory_path, "r", encoding="utf-8") as f:
            summarized_context = f.read()

def save_memory():
    memory_dir = f"C:\\Users\\{os.getlogin()}\\Anything AI\\memory\\"
    os.makedirs(memory_dir, exist_ok=True)
    with open(memory_dir + "memory.txt", "w", encoding="utf-8") as f:
        f.write(summarized_context)

def save_chat_history():
    if not context:
        return
    chat_dir = f"C:\\Users\\{os.getlogin()}\\Anything AI\\chats\\"
    os.makedirs(chat_dir, exist_ok=True)

    first_msg = ""
    for entry in context:
        if entry.startswith("User said:"):
            first_msg = entry.replace("User said:", "").strip()[:40]
            break

    safe_name = "".join(c for c in first_msg if c.isalnum() or c in (' ', '-', '_')).strip()
    safe_name = safe_name if safe_name else "chat"
    filename = safe_name + ".json"

    output = chat_dir + filename
    counter = 1
    while os.path.isfile(output):
        output = chat_dir + f"{safe_name}[{counter}].json"
        counter += 1

    with open(output, "w", encoding="utf-8") as f:
        json.dump(context, f, indent=2)

def get_recent_chats(n=5):
    chat_dir = f"C:\\Users\\{os.getlogin()}\\Anything AI\\chats\\"
    if not os.path.exists(chat_dir):
        return []
    files = sorted(
        [f for f in os.listdir(chat_dir) if f.endswith(".json")],
        reverse=True
    )
    return files[:n]

def load_chat(filename):
    global context
    chat_dir = f"C:\\Users\\{os.getlogin()}\\Anything AI\\chats\\"
    with open(chat_dir + filename, "r", encoding="utf-8") as f:
        context = json.load(f)

def print_chat_history():
    if not context:
        return
    print('\n')
    for entry in context:
        if entry.startswith('User said:'):
            print(col.WHITE + entry + col.WHITE)
        elif entry.startswith('Ai responded:'):
            print(col.LIGHTBLUE_EX + entry + col.WHITE)
        elif entry.startswith('AI opened') or entry.startswith('AI closed'):
            print(col.LIGHTGREEN_EX + entry + col.WHITE)
        elif entry.startswith('AI failed'):
            print(col.RED + entry + col.WHITE)
        elif entry.startswith('Ai asked to run'):
            print(col.YELLOW + entry + col.WHITE)
        else:
            print(col.WHITE + entry)
    print('\n')

# the most over the top, ANSI char spam code
def startup_menu():
    recent_chats = get_recent_chats()
    if not recent_chats:
        return None

    options = recent_chats + ["↾▰ Open Chats Folder..."] # most terminals cant even display the paraleelogram properly, it just looks like a smol square T_T
    selected = 0

    def draw_menu():
        for i, option in enumerate(options):
            name = option.replace(".json", "") if option.endswith(".json") else option
            prefix = col.LIGHTBLUE_EX + "⬡ " if i == selected else col.LIGHTBLACK_EX + "⬡ "
            print(f"\x1b[2K{prefix}{name}{col.WHITE}", flush=True)

    print(col.LIGHTBLACK_EX + "Ask Anything..." + col.WHITE)
    print()
    print(col.LIGHTBLACK_EX + "Recent Chats" + '=' * (width - len("Recent Chats")) + col.WHITE)
    print(f"\x1b[s", end="", flush=True)
    draw_menu()
    print(col.LIGHTBLACK_EX + '=' * width + col.WHITE + '\n')

    while True:
        ch = msvcrt.getwch()
        if ch == '\r':
            choice = options[selected]
            if choice.endswith(".json"):
                load_chat(choice)
                print_chat_history()
            else:
                os.startfile(f"C:\\Users\\{os.getlogin()}\\Anything AI\\chats\\")
            return None
        elif ch == '\xe0' or ch == '\x00':
            special = msvcrt.getwch()
            if special == 'H':
                selected = (selected - 1) % len(options)
                print(f"\x1b[u", end="", flush=True)
                draw_menu()
            elif special == 'P':
                selected = (selected + 1) % len(options)
                print(f"\x1b[u", end="", flush=True)
                draw_menu()
        elif ch.isprintable():
            os.system('cls')
            print('\n\n\n')
            ui_header()
            print('\n\n\n')
            return ch
        
def write_file_from_json(json_string): #big brain stuff
    try:
        data = json.loads(json_string)
    except json.JSONDecodeError as e:
        print(col.RED+f"JSON parse error: {e}")
        return

    filename = data.get("filename", "output.txt")
    title = data.get("title", "")
    content = data.get("content", [])
    ext = filename.split(".")[-1].lower()

    output_path = get_output_dir() + filename
    base, ext_part = os.path.splitext(filename)
    counter = 1
    while os.path.isfile(output_path):
        counter += 1
        filename = f"{base}[{counter}]{ext_part}"
        output_path = get_output_dir() + filename

    try:
        if ext == "pptx":
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN

            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            BLANK = prs.slide_layouts[6]

            def black_bg(slide):
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(0, 0, 0)

            def add_textbox(slide, text, left, top, width, height, font_size, bold=False, align=PP_ALIGN.LEFT):
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = text
                p.alignment = align
                run = p.runs[0]
                run.font.size = Pt(font_size)
                run.font.bold = bold
                run.font.color.rgb = RGBColor(255, 255, 255)
                return txBox

            title_slide = prs.slides.add_slide(BLANK)
            black_bg(title_slide)
            add_textbox(
                title_slide, title,
                left=Inches(1), top=Inches(2.8),
                width=Inches(11.33), height=Inches(1.5),
                font_size=48, bold=True, align=PP_ALIGN.CENTER
            )

            for slide_data in content:
                slide = prs.slides.add_slide(BLANK)
                black_bg(slide)

                add_textbox(
                    slide, slide_data.get("slide_title", ""),
                    left=Inches(0.5), top=Inches(0.3),
                    width=Inches(7.5), height=Inches(1),
                    font_size=28, bold=True
                )

                bullets = slide_data.get("bullets", [])

                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(7.5), Inches(5.5))
                tf = txBox.text_frame
                tf.word_wrap = True

                for i, b in enumerate(bullets):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = f"• {b}"
                    run = p.runs[0]
                    run.font.size = Pt(20)
                    run.font.color.rgb = RGBColor(255, 255, 255)

                img_query = slide_data.get("slide_title", "")
                img_stream = fetch_wikimedia_image(img_query)
                if img_stream:
                    try:
                        slide.shapes.add_picture(
                            img_stream,
                            left=Inches(8.3), top=Inches(1.5),
                            width=Inches(4.5), height=Inches(5.0)
                        )
                    except Exception:
                        pass

            prs.save(output_path)
            

        elif ext == "docx":
            doc = Document()
            doc.add_heading(title, 0)
            for section in content:
                doc.add_heading(section.get("section_title", ""), level=1)
                for para in section.get("paragraphs", []):
                    doc.add_paragraph(para)
            doc.save(output_path)

        elif ext == "pdf":
            c = canvas.Canvas(output_path, pagesize=A4)
            page_width, height = A4
            y = height - 50

            def check_newpage():
                nonlocal y
                if y < 60:
                    c.showPage()
                    y = height - 50

            c.setFont("Helvetica-Bold", 18)
            c.drawString(50, y, title)
            y -= 40

            for section in content:
                check_newpage()
                c.setFont("Helvetica-Bold", 13)
                c.drawString(50, y, section.get("section_title", ""))
                y -= 22

                c.setFont("Helvetica", 11)
                for para in section.get("paragraphs", []):
                    words = para.split()
                    line = ""
                    for word in words:
                        if c.stringWidth(line + word, "Helvetica", 11) < page_width - 100:
                            line += word + " "
                        else:
                            check_newpage()
                            c.drawString(50, y, line.strip())
                            y -= 18
                            line = word + " "
                    if line:
                        check_newpage()
                        c.drawString(50, y, line.strip())
                        y -= 18
                y -= 10

            c.save()

        elif ext in ("xlsx", "xls"):
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = title[:31]  # type: ignore
            for row in content:
                ws.append(row)  # type: ignore
            for cell in ws[1]:  # type: ignore
                cell.font = openpyxl.styles.Font(bold=True)  # type: ignore
            wb.save(output_path)

        elif ext == "csv":
            with open(output_path, "w", newline="", encoding="utf-8") as f:  # type: ignore
                writer = csv.writer(f)
                for row in content:
                    writer.writerow(row)

        elif ext in ("txt", "html", "md"):
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(title + "\n\n")
                for section in content:
                    f.write(section.get("section_title", "") + "\n")
                    for para in section.get("paragraphs", []):
                        f.write(para + "\n")
                    f.write("\n")

        type_print(col.LIGHTGREEN_EX+f"File saved to: {output_path}", 0.01)
        open_file(output_path)


    except Exception as e:
        print(col.RED+f"File write error: {e}")


def call_with_fallback(model, messages, max_tokens, temperature): # did this to cut 300 lines of my code, thats how bad it was
    for active_client in [client1, client2, client3]:
        try:
            response = active_client.chat.completions.create(
                model=model,
                messages=messages,
                max_tokens=max_tokens,
                temperature=temperature
            )
            return str(response.choices[0].message.content).strip()
        except:
            continue
    raise Exception("All available client tokens exhausted. Make sure your HF tokens are active.")


def sort_prompt(prompt):
    global context, summarized_context

    system_prompt = (
        "You are a central high-speed Intent Routing Node. Categorize the user prompt into exactly "
        "one of these strings: 'code_create', 'code_edit', 'reasoning', 'math', 'conversation', 'agentic_task', 'file_generate', 'terminal_command' or 'roleplay'.\n\n"
        
        "ROUTING RULES - read carefully:\n\n"
        
        "code_create: User wants you to WRITE source code for a program, script, or function. . DO NOT OUTPUT THIS IS THE USER IS ASKING ABOUT FILES"
        "Examples: 'write a snake game', 'make a python script that does X', 'code a sorting algorithm'.\n\n"
        
        "code_edit: User wants to modify, fix, debug, or improve existing code. "
        "ONLY use this if a code file is attached OR code exists in the conversation context below. "
        "Examples: 'fix this bug', 'add a feature to my script', 'refactor this function'.\n\n"
        
        "terminal_command: User wants to perform a FILE SYSTEM or OS-level operation that is best done with a single terminal command. "
        "Examples: 'create an empty text file', 'delete the folder', 'rename this file', 'list all files in directory', 'make a new folder'. "
        "If it can be done in one terminal command, route here instead of code_create.\n\n"
        
        "file_generate: User wants a document, spreadsheet, or presentation CREATED WITH CONTENT. DO NOT OUTPUT THIS IS THE USER IS ASKING ABOUT FILES"
        "Examples: 'make a powerpoint about space', 'create a word document explaining X', 'generate a csv of sales data'.\n\n"
        
        "agentic_task: User wants to OPEN or CLOSE an app or website on their PC. "
        "Examples: 'open spotify', 'close notepad', 'go to youtube.com'.\n\n"
        
        "reasoning: User has a complex logical, analytical, or diagnostic problem to solve step by step.\n\n"
        
        "math: User wants a mathematical calculation, proof, or algebraic manipulation.\n\n"
        
        "roleplay: User wants creative writing, storytelling, or character roleplay.\n\n"
        
        "conversation: Everything else — questions, explanations, general chat.\n\n"
        
        f"CONVERSATION CONTEXT (use to determine if code exists for code_edit):\n{get_context_string()}\n\n"
        
        "Output ONLY the exact category word. No punctuation, no explanation. Start immediately."
    )
    try:
        return call_with_fallback(
            "deepseek-ai/DeepSeek-R1-Distill-Llama-70B",
            [{"role": "system", "content": system_prompt}, {"role": "user", "content": str(prompt)}],
            50, 0.1
        )
    except:
        return "conversation"


def save_code_file(code): # save it? save it. save it, SAVE IT!
    lines = code.split("\n")
    suggested = "output.py"

    first_line = lines[0].strip()
    if "filename:" in first_line.lower():
        suggested = first_line.split("filename:")[-1].strip()
        code = "\n".join(lines[1:]).strip()

    savename = suggested
    type_print(f'Ok, saving the code as file {savename}', 0.01)
    output_path = get_output_dir() + savename
    if os.path.exists(output_path):
        type_print(f"'{suggested}' already exists. Overwrite? (y/n): ", 0.01)
        confirm = input().strip().lower()
        if confirm == "y":
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(code)
            type_print(f"File updated: {output_path}", 0.01)
            open_file(output_path)
        else:
            type_print("Save as new file instead? Enter filename or press Enter to discard: ", 0.01)
            newname = input().strip()
            if newname:
                output_path = get_output_dir() + newname
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(code)
                type_print(f"Saved as: {output_path}", 0.01)
                open_file(output_path)
            else:
                type_print("Discarded, file not saved.", 0.01)
    else:
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(code)
        type_print(f"Saved to {output_path}", 0.01)
        open_file(output_path)

# not really fits the definition of pipeline, it just sound hacky and cool 

def code_create_pipeline(prompt):
    global context, summarized_context

    if context:
        bg_thread = thread.Thread(target=maybe_summarize_context, daemon=True)
        bg_thread.start()

    system_prompt = (
        "You are Anything AI, an elite, hardware-level Senior Computer Science Architect and Compiler Engineer. "
        "Your task is to output syntactically flawless, execution-ready, and highly optimized code. "
        "Output ONLY raw source code. Do NOT use markdown code blocks or backticks (no ```). "
        "The very first line of your output must be a comment containing only the suggested filename. "
        "For python: '# filename: main.py', for js: '// filename: index.js', etc. Name what you see fit, for example, an insertion sort program may be called insertion_sort.py etc. "
        "Then start the actual code on line 2 immediately. "
        "Do not include any introductory text, comments, or explanations, just the bare code. Start writing the code immediately. "
        f"If necessary, use the conversation context given here: \n{get_context_string()}"
    )

    try:
        ai_reply = call_with_fallback(
            "Qwen/Qwen3-Coder-480B-A35B-Instruct",
            [{"role": "system", "content": system_prompt}, {"role": "user", "content": str(prompt)}],
            2000, 0.3
        )
        context.append('User said: ' + str(prompt))
        context.append('Ai responded: ' + ai_reply)
        return ai_reply
    except Exception as e:
        return f"Native Hub Error: {e}\nMake sure your HF token is valid and active."


def code_edit_pipeline(prompt):
    global context, summarized_context

    if context:
        bg_thread = thread.Thread(target=maybe_summarize_context, daemon=True)
        bg_thread.start()

    system_prompt = (
        "You are Anything AI, an elite, world-class Senior Systems Engineer and Debugging Expert. Your task is to analyze code, find bugs, and apply modifications based on user instructions."
        "\n\nCRITICAL CONTEXT RULES:\nBelow is the history of this conversation session. Use it to understand the evolution of the script, but prioritize the user's latest prompt:"
        f"\n{get_context_string()}\nOutput ONLY the modified, execution-ready code. Do not use markdown blocks, backticks, or conversational commentary, but do put comments in the code itself using that language's correct comment syntax detailing what you have changed. "
        "There is probably a filename comment at the very start of the code you are editing, make sure to NOT change the filename at ALL"
        "Start the fixed code immediately."
    )

    try:
        ai_reply = call_with_fallback(
            "Qwen/Qwen3-Coder-480B-A35B-Instruct",
            [{"role": "system", "content": system_prompt}, {"role": "user", "content": str(prompt)}],
            2000, 0.1
        )
        context.append('User said: ' + str(prompt))
        context.append('Ai responded: ' + ai_reply)
        return ai_reply
    except Exception as e:
        return f"Native Hub Error: {e}\nMake sure your HF token is valid and active."


def conversation_pipeline(prompt):
    global context, summarized_context

    if context:
        bg_thread = thread.Thread(target=maybe_summarize_context, daemon=True)
        bg_thread.start()

    system_prompt = (
        "You are Anything AI, a helpful, brilliant, and slightly tech-witty AI desktop companion. Provide insightful, grounded, and concise conversational responses. Make sure your response has no markdown formatting and is in plaintext, something can be pasted in a terminal window.\n\n"
        f"Use this context to keep track of the discussion flow:\n{get_context_string()}"
    )

    try:
        ai_reply = call_with_fallback(
            "Qwen/Qwen2.5-72B-Instruct",
            [{"role": "system", "content": system_prompt}, {"role": "user", "content": str(prompt)}],
            1000, 0.7
        )
        context.append('User said: ' + str(prompt))
        context.append('Ai responded: ' + ai_reply)
        return ai_reply
    except Exception as e:
        return f"Native Hub Error: {e}\nMake sure your HF token is valid and active."


def roleplay_pipeline(prompt):
    global context, summarized_context

    if context:
        bg_thread = thread.Thread(target=maybe_summarize_context, daemon=True)
        bg_thread.start()

    system_prompt = (
        "You are Anything AI, an immersive, highly adaptive creative writing and roleplay engine. Maintain strict adherence to characters, settings, personas, and tones established by the user. Stay in character completely. Make sure your response has no markdown formatting and is in plaintext, something can be pasted in a terminal window.\n\n"
        f"Prior narrative context:\n{get_context_string()}"
    )

    try:
        ai_reply = call_with_fallback(
            "Qwen/Qwen2.5-72B-Instruct",
            [{"role": "system", "content": system_prompt}, {"role": "user", "content": str(prompt)}],
            1200, 0.85
        )
        context.append('User said: ' + str(prompt))
        context.append('Ai responded: ' + ai_reply)
        return ai_reply
    except Exception as e:
        return f"Native Hub Error: {e}\nMake sure your HF token is valid and active."


def reasoning_pipeline(prompt):
    global context, summarized_context

    if context:
        bg_thread = thread.Thread(target=maybe_summarize_context, daemon=True)
        bg_thread.start()

    system_prompt = (
        "You are Anything AI, an advanced logical inference core. Break down the user's complex query, puzzle, or diagnostic trace step-by-step with absolute accuracy. Make sure your response has no markdown formatting and is in plaintext, something can be pasted in a terminal window, but should still contain puctuation like capitalization full stops, commas, colons, etc..\n\n"
        f"Relevant session data:\n{get_context_string()}"
    )

    try:
        ai_reply = call_with_fallback(
            "deepseek-ai/DeepSeek-R1-Distill-Llama-70B",
            [{"role": "system", "content": system_prompt}, {"role": "user", "content": str(prompt)}],
            2500, 0.2
        )
        context.append('User said: ' + str(prompt))
        context.append('Ai responded: ' + ai_reply)
        return ai_reply
    except Exception as e:
        return f"Native Hub Error: {e}\nMake sure your HF token is valid and active."


def math_pipeline(prompt):
    global context, summarized_context

    if context:
        bg_thread = thread.Thread(target=maybe_summarize_context, daemon=True)
        bg_thread.start()

    system_prompt = (
        "You are Anything AI, a specialized mathematical calculation and verification engine. Perform precise computation, algebraic manipulation, or proof validation step-by-step. Make sure your response has no markdown formatting and is in plaintext, something can be pasted in a terminal window, but should still contain puctuation like capitalization full stops, commas, colons, etc.\n\n"
        f"Mathematical context history:\n{get_context_string()}"
    )

    try:
        ai_reply = call_with_fallback(
            "deepseek-ai/DeepSeek-R1-Distill-Llama-70B",
            [{"role": "system", "content": system_prompt}, {"role": "user", "content": str(prompt)}],
            1500, 0.1
        )
        context.append('User said: ' + str(prompt))
        context.append('Ai responded: ' + ai_reply)
        return ai_reply
    except Exception as e:
        return f"Native Hub Error: {e}\nMake sure your HF token is valid and active."


def agentic_task_pipeline(prompt):
    global context, summarized_context

    if context:
        bg_thread = thread.Thread(target=maybe_summarize_context, daemon=True)
        bg_thread.start()

    applist = app.give_appnames()
    system_prompt = (
        "You are Anything AI, a local workstation automation execution bridge. Your sole job is to translate the user's request into a single system command string.\n\n"
        "CRITICAL RULES:\n"
        "1. If the user wants to launch or start an app, output EXACTLY 'app:open:[app name]'. Example: app:open:notepad\n"
        "2. If the user wants to terminate or close an app, output EXACTLY 'app:close:[app name]'. Example: app:close:notepad\n"
        "3. If the user wants to view a website or internet link, output EXACTLY 'web:[website url]'. Example: web:https://google.com\n"
        "4. Output ONLY the exact command string format. Do not use markdown blocks, backticks, or trailing punctuation. Do not explain anything.\n\n"
        f"AVAILABLE SYSTEM APPLICATION LIST:\n"
        f"You can ONLY target applications present in this list when generating 'app:' commands:\n[{applist}]\n\n"
        "5. DESCRIPTIVE CLAUSE: If the user describes an application or website by its purpose instead of its direct title, you must deduce the proper tool. "
        "If they describe an application, select the closest functional match available in the application list above.\n\n"
        "6. SPECIFIC WEBSITE PAGE CLAUSE:\n"
        "If the user also describes or names a certain page of the website, open that specific page, for example, if a user asks to open the My Stuff page on scratch, return the URL : scratch.mit.edu/mystuff"
        f"Session context:\n{get_context_string()}"
    )

    def dispatch(ai_reply):
        if ai_reply.startswith('app:open:'):
            appname = ai_reply.removeprefix('app:open:')
            try:
                app.open(appname, match_closest=True, output=False, throw_error=True)
                type_print(col.LIGHTGREEN_EX + 'AI opened app: ' + appname, 0.01)
                context.append('AI opened app: ' + appname)
            except Exception as e:
                type_print(col.LIGHTBLUE_EX + "Sorry, I couldn't find that application on your system.", 0.01)
                context.append('AI failed to open app: ' + appname)
        elif ai_reply.startswith('app:close:'):
            appname = ai_reply.removeprefix('app:close:')
            try:
                app.close(appname, match_closest=True, output=False, throw_error=True)
                type_print(col.LIGHTGREEN_EX + 'AI closed app: ' + appname, 0.01)
                context.append('AI closed app: ' + appname)
            except Exception as e:
                type_print(col.LIGHTBLUE_EX + "Sorry, I couldn't find that application running on your system.", 0.01)
                context.append('AI failed to close app: ' + appname)
        elif ai_reply.startswith('web:'):
            webname = ai_reply.removeprefix('web:')
            try:
                web.open(webname)
                type_print(col.LIGHTGREEN_EX + 'AI opened website: ' + webname, 0.01)
                context.append('AI opened website: ' + webname)
            except Exception as e:
                type_print(col.LIGHTBLUE_EX + "Sorry, I couldn't open the website you specified", 0.01)
                context.append('AI failed to open website: ' + webname)

    try:
        ai_reply = call_with_fallback(
            "Qwen/Qwen2.5-Coder-32B-Instruct",
            [{"role": "system", "content": system_prompt}, {"role": "user", "content": str(prompt)}],
            50, 0.1
        )
        context.append('User said: ' + str(prompt))
        dispatch(ai_reply)
    except Exception as e:
        return f"Native Hub Error: {e}\nMake sure your HF token is valid and active."
            
def file_generate_pipeline(prompt):
    global context, summarized_context

    if context:
        bg_thread = thread.Thread(target=maybe_summarize_context, daemon=True)
        bg_thread.start()

    system_prompt = (
        "You are Anything AI, a structured data compilation core. Your operational purpose is to convert human requests for documents, sheets, or presentations into a strictly compliant, standardized JSON schema.\n\n"
        "JSON SCHEMA STRUCTURE:\n"
        "Your output must strictly follow this exact structural signature:\n"
        "{\n"
        '  "filename": "requested_name.extension",\n'
        '  "title": "Main Document Title or Header",\n'
        '  "content": []\n'
        "}\n\n"
        "DATA HANDLING RULES:\n"
        "1. For text documents (.docx, .pdf, .html, .txt): The 'content' array must be a list of strings where each element represents a paragraph, bullet point, or section body.\n"
        "2. For presentation outlines (.pptx, .ppt): The 'content' array must be a list of strings where each element represents a slide's key bullet points or structural notes.\n"
        "3. For spreadsheets and data tables (.xlsx, .xls, .csv): The 'content' array must be a two-dimensional array (an array of arrays), where each inner array represents a horizontal row of cells. The first row should contain the headers. Example: [[\"Header1\", \"Header2\"], [\"Row1Cell1\", \"Row1Cell2\"]]\n\n"
        "CRITICAL EXECUTION CONSTRAINTS:\n"
        "- Do NOT wrap your output in markdown formatting like ```json ... ```.\n"
        "- Do NOT include any introductory or concluding text, warnings, or conversational pleasantries.\n"
        "- Start your response immediately with the opening curly brace '{' and end immediately with the closing curly brace '}'.\n"
        "- Ensure all inner quotation marks inside the strings are properly escaped to prevent parsing errors.\n"
        "- Start working immediately."
        "For presentation files (.pptx): The 'content' array must be a list of slide objects, each with this exact structure:\n"
        "{ \"slide_title\": \"Title of the slide\", \"bullets\": [\"bullet 1\", \"bullet 2\", \"bullet 3\"] }\n"
        "Example: [{\"slide_title\": \"Introduction\", \"bullets\": [\"Point one\", \"Point two\"]}, ...]\n"
        "Do NOT use flat strings for slides. Each slide MUST be a JSON object with 'slide_title' and 'bullets' keys.\n"
        "For text documents (.docx, .pdf, .txt): The 'content' array must be a list of section objects:\n"
        "{ \"section_title\": \"Section Heading\", \"paragraphs\": [\"paragraph 1\", \"paragraph 2\"] }\n"
        "Example: [{\"section_title\": \"Introduction\", \"paragraphs\": [\"This is the intro.\", \"More detail here.\"]}]\n\n"
        "For spreadsheets (.xlsx, .csv): The 'content' array must be a two-dimensional array where the first row is headers:\n"
        "Example: [[\"Name\", \"Age\", \"City\"], [\"Alice\", \"30\", \"Mumbai\"], [\"Bob\", \"25\", \"Delhi\"]]\n"
        "Each inner array is one row. All values must be strings.\n\n"
        "For PDF (.pdf): Use the SAME section object format as .docx:\n"
        "{ \"section_title\": \"Section Heading\", \"paragraphs\": [\"paragraph 1\", \"paragraph 2\"] }\n"
    )
    try:
        ai_reply = call_with_fallback(
            "Qwen/Qwen3-Coder-480B-A35B-Instruct",
            [{"role": "system", "content": system_prompt}, {"role": "user", "content": str(prompt)}],
            2000, 0.1
        )
        context.append('User said: ' + str(prompt))
        context.append('Ai responded: ' + ai_reply)
        return ai_reply
    except Exception as e:
        return f"Native Hub Error: {e}\nMake sure your HF token is valid and active."

def image_pipeline(image_path, user_question="Describe this image deeply."):
    global context, summarized_context

    # Convert the local image to a base64 Data URI string
    data_url = get_image_base64_uri(image_path)

    system_prompt = (
        "You are Anything AI, an advanced visual preprocessing engine. Analyze the provided image deeply. "
        "Provide a highly detailed, text-only breakdown describing the contents, user interface elements, "
        "code snippets, text, errors, or objects visible. Your description will be read by another AI model, "
        "so be precise, objective, and dense with information. Do not use any markdown formatting.\n"
        f"If necessary, use the conversation context given here: \n{get_context_string()}"
    )

    clients = [client1, client2, client3]
    
    messages_payload = [
        {"role": "system", "content": system_prompt},
        {
            "role": "user",
            "content": [# hello, open source enthusiast, if u reached here reading my sphaggeti code, congrats!
                {"type": "text", "text": str(user_question)},
                {"type": "image_url", "image_url": {"url": data_url}}
            ]
        }
    ]

    for idx, active_client in enumerate(clients):
        try: 
            response = active_client.chat.completions.create(
                model="Qwen/Qwen3.5-397B-A17B",
                messages=messages_payload,
                max_tokens=1000,
                temperature=0.1
            )
            
            ai_reply = str(response.choices[0].message.content).strip()
            
            context.append(f"User asked about image ({os.path.basename(image_path)}): {user_question}")
            context.append(f"Ai responded: {ai_reply}")
            return ai_reply
            
        except Exception as e:
            continue
    return "Native Hub Error: All available client tokens exhausted. Make sure your HF tokens are active."

def terminal_command_pipeline(prompt):
    global context, summarized_context
    current_os = platform.system()
    if context:
        bg_thread = thread.Thread(target=maybe_summarize_context, daemon=True)
        bg_thread.start()

    system_prompt = (
        f"You are a text-to-CLI translation engine. You must convert human intent into a functional terminal command for the {current_os} operating system.\n"
        "ALLOWED OUTPUT FORMAT:\n"
        "mkdir test_folder\n\n"
        "FORBIDDEN OUTPUT FORMATS:\n"
        "- ```bash\\nmkdir test_folder\\n```\n"
        "- Here is the command: mkdir test_folder\n"
        "- mkdir test_folder (Note: this creates a directory)\n\n"
        "Strictly evaluate inputs and output the raw plain-text command only. No introductions. No conclusions."
    )
    try:
        ai_reply = call_with_fallback(
            "Qwen/Qwen3-Coder-480B-A35B-Instruct",
            [{"role": "system", "content": system_prompt}, {"role": "user", "content": str(prompt)}],
            2000, 0.1
        ).replace("`", "")
        context.append('User said: ' + str(prompt))
        context.append('Ai asked to run command: ' + ai_reply)
        return ai_reply
    except Exception as e:
        return f"Native Hub Error: {e}\nMake sure your HF token is valid and active." 


filebar_files = []

def draw_file_bar():
    if filebar_files:
        bar = "  ".join(f"[ {f[0]} {f[1]} ]" for f in filebar_files)
        print(f"\x1b[u\x1b[2K{col.LIGHTGREEN_EX}{bar}{col.WHITE}\x1b[u\x1b[2B", end="", flush=True)



if __name__ == "__main__": # first time ive evr used if __name__ == __main__ btw
    attachfile = []
    print('\n\n\n')
    ui_header()
    print('\n\n\n')
    load_memory()
    seed = startup_menu()
    print(col.LIGHTBLACK_EX + (' ' * (width - 26)) + '( + [Insert] | → [Enter] )\n\n')

    while True:
        filebar_files = []

        print(f"\x1b[s")
        print(col.LIGHTGREEN_EX + '' + col.WHITE)

        task = read_input(col.LIGHTBLACK_EX + "Ask Anything..." + col.WHITE, seed=seed)
        seed = None 
        if task.strip() == "":
            if context:
                print_chat_history()
            continue

        print()
        print(col.LIGHTBLACK_EX + 'Processing Files...\r', end='', flush=True)
        for every_file in attachfile:
            fileext = every_file.split('.')[-1].lower()

            if fileext in ("png", "jpg", "jpeg", "gif", "webp"):
                texted_image = image_pipeline(every_file)
                task = task + ', User also attached image file showing: ' + texted_image
            elif fileext == 'txt':
                parsedfile = parse_text_file(every_file)
                task = task + ", User also attached text file with contents:\n" + parsedfile
            elif fileext == 'csv':
                parsedfile = parse_csv_file(every_file)
                task = task + ", User also attached csv file with contents:\n" + parsedfile
            elif fileext == 'docx':
                parsedfile = parse_docx_file(every_file)
                task = task + ", User also attached Microsoft Word file with contents:\n" + parsedfile
            elif fileext == 'xlsx':
                parsedfile = parse_xlsx_file(every_file)
                task = task + ", User also attached Microsoft Excel file with contents:\n" + parsedfile
            elif fileext == 'pptx':
                parsedfile = parse_pptx_file(every_file)
                task = task + ", User also attached Microsoft PowerPoint file with contents:\n" + parsedfile
            elif fileext in ("py", "java", "js", "html", "kt", "c", "cpp", "cs", "css", "pyw", "ipynb"):
                parsedfile = parse_code_file(every_file)
                task = task + f", User also attached {fileext} code file with contents:\n" + parsedfile
            else:
                parsedfile = parse_text_file(every_file)
                task = task + ", User also attached file with contents:\n" + parsedfile

        attachfile = []
        print('                              ')
        intent = str(sort_prompt(task)).lower()

        if intent == 'code_create':
            ai_response = code_create_pipeline(task)
            type_print(col.LIGHTBLUE_EX + str(ai_response), 0.01)
            type_print(col.LIGHTYELLOW_EX + '\nWould you like me to save this file to your computer? [y/n]', 0.01)
            chocice = str(input())
            while chocice.lower() not in ('y', 'n'):
                type_print(col.LIGHTYELLOW_EX + 'Please enter y or n: ', 0.01)
                chocice = str(input())
            if chocice.lower() == 'y':
                save_code_file(ai_response)
            else:
                continue
        elif intent == 'code_edit':
            ai_response = code_edit_pipeline(task)
            type_print(col.LIGHTBLUE_EX + str(ai_response), 0.01)
            type_print(col.LIGHTYELLOW_EX + '\nWould you like me to save/edit this file on your computer? [y/n]', 0.01)
            chocice = str(input())
            while chocice.lower() not in ('y', 'n'):
                type_print(col.LIGHTYELLOW_EX + 'Please enter y or n: ', 0.01)
                chocice = str(input())
            if chocice.lower() == 'y':
                save_code_file(ai_response)
            else:
                continue
        elif intent == 'reasoning':
            ai_response = reasoning_pipeline(task)
            type_print(col.LIGHTBLUE_EX + str(ai_response), 0.01)
        elif intent == 'conversation':
            ai_response = conversation_pipeline(task)
            type_print(col.LIGHTBLUE_EX + str(ai_response), 0.01)
        elif intent == 'math':
            ai_response = math_pipeline(task)
            type_print(col.LIGHTBLUE_EX + str(ai_response), 0.01)
        elif intent == 'roleplay':
            ai_response = roleplay_pipeline(task)
            type_print(col.LIGHTBLUE_EX + str(ai_response), 0.01)
        elif intent == 'agentic_task':
            agentic_task_pipeline(task)
            ai_response = ""
        elif intent == 'file_generate':
            ai_response = file_generate_pipeline(task)
            write_file_from_json(ai_response)
        elif intent == 'terminal_command':
            run_terminal_command(terminal_command_pipeline(task))
        else:
            ai_response = conversation_pipeline(task)
            type_print(col.LIGHTBLUE_EX + str(ai_response), 0.01)

        save_memory()
        save_chat_history()

# namaskar sahab, aapne itna niche kaise scroll kiya
